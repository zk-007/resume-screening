import os
import pandas as pd

def read_pdf(file) -> str:
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                text_parts.append(page.extract_text() or "")
        return "\n".join(text_parts).strip()
    except Exception:
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(file)
            pages = [p.extract_text() or "" for p in reader.pages]
            return "\n".join(pages).strip()
        except Exception:
            return ""

def read_docx(file) -> str:
    try:
        import docx2txt
        bytes_data = file.read()
        temp_path = "tmp_upload.docx"
        with open(temp_path, "wb") as f:
            f.write(bytes_data)
        text = docx2txt.process(temp_path) or ""
        os.remove(temp_path)
        return text.strip()
    except Exception:
        return ""

def read_rtf(file) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        data = file.read().decode(errors="ignore")
        return rtf_to_text(data).strip()
    except Exception:
        return ""

def read_txt(file) -> str:
    try:
        return file.read().decode(errors="ignore").strip()
    except Exception:
        try:
            return file.read().decode("utf-8", errors="ignore").strip()
        except Exception:
            return ""

def read_excel(file) -> str:
    try:
        df = pd.read_excel(file)
        return "\n".join(df.astype(str).fillna("").values.ravel().tolist()).strip()
    except Exception:
        return ""

def read_pptx(file) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts).strip()
    except Exception:
        return ""

READERS = {
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".rtf": read_rtf,
    ".txt": read_txt,
    ".xlsx": read_excel,
    ".xls": read_excel,
    ".pptx": read_pptx,
}